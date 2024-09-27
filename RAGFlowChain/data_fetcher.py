from bs4 import BeautifulSoup
import requests
from langchain_community.document_loaders import PyPDFLoader, DataFrameLoader, YoutubeLoader, TextLoader
from langchain.text_splitter import CharacterTextSplitter
from langchain.schema import Document
from docx import Document as DocxDocument
from pptx import Presentation as PptxPresentation
from tempfile import NamedTemporaryFile
from googleapiclient.discovery import build
import pandas as pd
import os
import yaml
from youtube_transcript_api import YouTubeTranscriptApi, NoTranscriptFound

# Function to extract text from a webpage
def extract_text_from_webpage(url):
    headers = {
        "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36"
    }
    response = requests.get(url, headers=headers)
    
    # Check for successful response
    if response.status_code == 200:
        soup = BeautifulSoup(response.content, 'html.parser')
        text = soup.get_text(separator=' ', strip=True)
        return text
    else:
        return f"Error {response.status_code}: Unable to access {url}"

# Function to fetch website content
def fetch_website_content(url, chunk_size=1000):
    try:
        content = extract_text_from_webpage(url)
        
        # Handle error messages returned by extract_text_from_webpage
        if content.startswith("Error"):
            print(content)
            return []
        
        # Split the website content into manageable chunks
        text_splitter = CharacterTextSplitter(chunk_size=chunk_size, separator="\n")
        documents = text_splitter.split_text(content)
        
        # Convert to list of dicts, similar to other sources
        document_data = [{"source": url, "content": chunk, "source_type": "website"} for chunk in documents]
        return document_data
    
    except requests.exceptions.RequestException as e:
        print(f"Error fetching website content from {url}: {e}")
        return []

def load_and_split_ppt(ppt_path, chunk_size):
    # Load the PowerPoint file
    presentation = Presentation(ppt_path)
    
    # Extract text from all slides
    full_text = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                full_text.append(shape.text)
    
    # Combine all text and split into chunks
    full_text = "\n".join(full_text)
    text_splitter = CharacterTextSplitter(chunk_size=chunk_size, separator="\n")
    documents = text_splitter.split_text(full_text)
    
    # Convert to list of dicts for consistency with other data sources
    document_data = [{"source": ppt_path, "content": chunk, "source_type": "ppt"} for chunk in documents]
    return document_data

# Function to load and split documents from local paths
def load_and_split_documents(file_or_dir_path, chunk_size=1000):
    all_documents = []

    if os.path.isdir(file_or_dir_path):
        for root, _, files in os.walk(file_or_dir_path):
            for file in files:
                file_path = os.path.join(root, file)
                all_documents.extend(load_and_split_documents(file_path, chunk_size))
    else:
        _, file_extension = os.path.splitext(file_or_dir_path)

        if file_extension == ".pdf":
            with NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
                with open(file_or_dir_path, 'rb') as src:
                    tmp.write(src.read())
                file_path = tmp.name
            loader = PyPDFLoader(file_path)
            docs = loader.load()

        elif file_extension == ".pptx":
            prs = PptxPresentation(file_or_dir_path)
            docs = [
                Document(
                    page_content="\n".join(
                        [shape.text for shape in slide.shapes if hasattr(shape, "text")]
                    ),
                    metadata={"title": slide.shapes.title.text if hasattr(slide.shapes.title, "text") else ""}
                )
                for slide in prs.slides
            ]

        elif file_extension == ".docx":
            doc = DocxDocument(file_or_dir_path)
            docs = [
                Document(
                    page_content=para.text,
                    metadata={}
                )
                for para in doc.paragraphs
            ]

        elif file_extension == ".txt":
            with open(file_or_dir_path, 'r', encoding='utf-8') as txt_file:
                full_text = txt_file.read()
            docs = [Document(page_content=full_text, metadata={})]

        else:
            print(f"Skipping unsupported file type: {file_extension}")
            return []

        text_splitter = CharacterTextSplitter(chunk_size=chunk_size, separator="\n")
        split_docs = text_splitter.split_documents(docs)

        document_data = [
            {"source": file_or_dir_path, 
             "content": doc.page_content.replace('\n\n', '\n'),  # Replace double newlines with a single newline
             "source_type": file_extension[1:], 
             "metadata": doc.metadata}
            for doc in split_docs
        ]

        all_documents.extend(document_data)

    return all_documents

# Function to fetch books from Google Books API
def fetch_books(api_key, query, max_results=10):
    service = build('books', 'v1', developerKey=api_key)
    request = service.volumes().list(q=query, maxResults=max_results)
    response = request.execute()
    books = response.get('items', [])
    return books

def extract_books_data(books):
    data = []
    for book in books:
        volume_info = book.get('volumeInfo', {})
        book_data = {
            "source": book.get('id'),
            "title": volume_info.get('title'),
            "author": ", ".join(volume_info.get('authors', [])),
            "publishedDate": volume_info.get('publishedDate'),
            "description": volume_info.get('description', ''),
            "content": volume_info.get('description', '').replace('\n\n', '\n')  # Clean content
        }
        data.append(book_data)
    return data

# Function to fetch news articles from NewsAPI
def fetch_news_articles(api_key, query, page_size=5, max_pages=1):
    all_articles = []
    for page in range(1, max_pages + 1):
        url = "https://newsapi.org/v2/everything"
        params = {
            'apiKey': api_key,
            'q': query,
            'pageSize': page_size,
            'page': page
        }
        response = requests.get(url, params=params)
        if response.status_code == 426:
            print("Upgrade required to access more data.")
            break
        response.raise_for_status()
        articles = response.json().get('articles', [])
        if not articles:
            break
        all_articles.extend(articles)
    return all_articles

def extract_articles_data(articles):
    data = []
    for article in articles:
        article_data = {
            "source": article['source']['name'],
            "author": article.get('author', ''),
            "title": article.get('title', ''),
            "description": article.get('description', ''),
            "url": article.get('url', ''),
            "publishedDate": article.get('publishedAt', ''),
            "content": article.get('content', '').replace('\n\n', '\n')  # Clean content
        }
        data.append(article_data)
    return data

# Function to search for videos on YouTube
def search_videos(topic, api_key, max_results=10):
    youtube = build('youtube', 'v3', developerKey=api_key)
    request = youtube.search().list(
        q=topic,
        part='id,snippet',
        maxResults=max_results,
        type='video'
    )
    response = request.execute()
    video_ids = [item['id']['videoId'] for item in response['items']]
    return video_ids

def load_video(video_id):
    url = f'https://www.youtube.com/watch?v={video_id}'
    loader = YoutubeLoader.from_youtube_url(url, add_video_info=True)
    try:
        docs = loader.load()
        if not docs:
            return None
        doc = docs[0]
        doc_df = pd.DataFrame([doc.metadata])
        doc_df['url'] = url
        doc_df['content'] = doc.page_content.replace('\n\n', '\n')  # Clean content
        return doc_df
    except NoTranscriptFound:
        return None

# Function to fetch data from multiple sources
def fetch_data(online_sources=None, local_sources=None, chunk_size=1000):
    all_data = []

    if online_sources:
        youtube = online_sources.get('youtube')
        websites = online_sources.get('websites')
        books = online_sources.get('books')
        news_articles = online_sources.get('news_articles')
        
        if youtube:
            youtube_query = youtube.get('topic')
            youtube_api_key = youtube.get('api_key')
            youtube_max_results = youtube.get('max_results')
            if youtube_query and youtube_api_key and youtube_max_results:
                video_ids = search_videos(youtube_query, youtube_api_key, youtube_max_results)
                videos = []
                for video_id in video_ids:
                    video = load_video(video_id)
                    if video is not None:
                        video['content'] = video['content'].str.replace('\n\n', '\n', regex=False)  # Replace double newlines
                        videos.append(video)
                
                if videos:
                    youtube_data = pd.concat(videos, ignore_index=True)
                    youtube_data['source_type'] = 'youtube'
                    all_data.extend(youtube_data.to_dict('records'))

        if websites:
            for url in websites:
                website_data = fetch_website_content(url, chunk_size=chunk_size)
                for data in website_data:
                    data['content'] = data['content'].replace('\n\n', '\n')  # Replace double newlines
                all_data.extend(website_data)

        if books:
            books_api_key = books.get('api_key')
            books_query = books.get('query')
            books_max_results = books.get('max_results')
            if books_api_key and books_query and books_max_results:
                books_data = fetch_books(books_api_key, books_query, books_max_results)
                if books_data:
                    extracted_books_data = extract_books_data(books_data)
                    for item in extracted_books_data:
                        item['content'] = item['content'].replace('\n\n', '\n')  # Replace double newlines
                        item['source_type'] = 'book'
                    all_data.extend(extracted_books_data)

        if news_articles:
            news_api_key = news_articles.get('api_key')
            news_query = news_articles.get('query')
            news_page_size = news_articles.get('page_size')
            news_max_pages = news_articles.get('max_pages')
            if news_api_key and news_query and news_page_size and news_max_pages:
                news_articles_data = fetch_news_articles(news_api_key, news_query, news_page_size, news_max_pages)
                if news_articles_data:
                    extracted_news_data = extract_articles_data(news_articles_data)
                    for item in extracted_news_data:
                        item['content'] = item['content'].replace('\n\n', '\n') 
                        item['source_type'] = 'news'
                    all_data.extend(extracted_news_data)

    if local_sources:
        for local_source in local_sources:
            document_data = load_and_split_documents(local_source, chunk_size=chunk_size)
            all_data.extend(document_data)

    if all_data:
        final_data = pd.DataFrame(all_data)
        standardized_columns = ['source', 'title', 'author', 'publishedDate', 'description', 'content', 'url', 'source_type']
        for col in standardized_columns:
            if col not in final_data.columns:
                final_data[col] = None
        final_data = final_data[standardized_columns]
    else:
        final_data = pd.DataFrame()

    return final_data


